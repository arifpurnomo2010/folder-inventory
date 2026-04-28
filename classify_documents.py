"""
classify_documents.py
─────────────────────
Baca ratusan dokumen lokal (PDF, DOCX, XLSX, PPTX) → identifikasi industri via Claude API → output CSV.

SETUP:
    pip install anthropic pypdf python-docx openpyxl python-pptx tqdm

CARA PAKAI:
    1. Set ANTHROPIC_API_KEY di environment variable, atau isi langsung di baris ANTHROPIC_API_KEY di bawah.
    2. Ubah FOLDER_PATH ke folder dokumen kamu.
    3. Jalankan: python classify_documents.py
    4. Hasil tersimpan di: hasil_klasifikasi.csv
"""

import os
import csv
import time
import traceback
from pathlib import Path

# ── Install check ────────────────────────────────────────────────────────────
try:
    import anthropic
    from pypdf import PdfReader
    import docx
    import openpyxl
    from pptx import Presentation
    from tqdm import tqdm
except ImportError as e:
    print(f"[ERROR] Library belum terinstall: {e}")
    print("Jalankan: pip install anthropic pypdf python-docx openpyxl python-pptx tqdm")
    exit(1)

# ── KONFIGURASI — ubah sesuai kebutuhan ──────────────────────────────────────
FOLDER_PATH       = r"C:\Users\NamaKamu\Documents\Dokumen"   # ← ganti ini
OUTPUT_CSV        = "hasil_klasifikasi.csv"
ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")  # atau isi langsung: "sk-ant-..."

MAX_CHARS_PER_DOC = 3000   # Cukup untuk klasifikasi, hemat token
DELAY_BETWEEN_API = 0.5    # Detik jeda antar API call (hindari rate limit)
REKURSI_SUBFOLDER = True   # True = masuk subfolder juga

KATEGORI_INDUSTRI = [
    "Susu & Dairy", "Snack & Makanan Ringan", "Minuman", "Makanan & F&B",
    "Shampo & Perawatan Rambut", "Skincare & Kecantikan", "Kesehatan & Farmasi",
    "Otomotif (Motor)", "Otomotif (Mobil)", "Perbankan & Keuangan",
    "Asuransi", "Telekomunikasi", "E-commerce & Retail", "Properti",
    "Pendidikan", "Logistik & Transportasi", "Media & Hiburan",
    "Teknologi & Software", "Manufaktur", "Pemerintahan",
    "Lainnya (tidak teridentifikasi)"
]

# ─────────────────────────────────────────────────────────────────────────────

def ekstrak_teks_pdf(path: Path) -> str:
    try:
        reader = PdfReader(str(path))
        teks = ""
        for page in reader.pages:
            teks += page.extract_text() or ""
            if len(teks) >= MAX_CHARS_PER_DOC:
                break
        return teks[:MAX_CHARS_PER_DOC]
    except Exception:
        return ""


def ekstrak_teks_docx(path: Path) -> str:
    try:
        doc = docx.Document(str(path))
        teks = "\n".join(p.text for p in doc.paragraphs)
        return teks[:MAX_CHARS_PER_DOC]
    except Exception:
        return ""


def ekstrak_teks_xlsx(path: Path) -> str:
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        teks = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            teks += f"[Sheet: {sheet}]\n"
            for row in ws.iter_rows(values_only=True):
                baris = " | ".join(str(c) for c in row if c is not None)
                if baris:
                    teks += baris + "\n"
                if len(teks) >= MAX_CHARS_PER_DOC:
                    break
            if len(teks) >= MAX_CHARS_PER_DOC:
                break
        return teks[:MAX_CHARS_PER_DOC]
    except Exception:
        return ""


def ekstrak_teks_pptx(path: Path) -> str:
    try:
        prs = Presentation(str(path))
        teks = ""
        for i, slide in enumerate(prs.slides):
            teks += f"[Slide {i+1}]\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        teks += para.text + "\n"
            if len(teks) >= MAX_CHARS_PER_DOC:
                break
        return teks[:MAX_CHARS_PER_DOC]
    except Exception:
        return ""


def ekstrak_teks(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return ekstrak_teks_pdf(path)
    elif ext == ".docx":
        return ekstrak_teks_docx(path)
    elif ext in (".xlsx", ".xlsm"):
        return ekstrak_teks_xlsx(path)
    elif ext == ".pptx":
        return ekstrak_teks_pptx(path)
    return ""


def klasifikasi_industri(client: anthropic.Anthropic, filename: str, teks: str) -> dict:
    if not teks.strip():
        return {
            "industri": "Lainnya (tidak teridentifikasi)",
            "keyakinan": "rendah",
            "alasan": "Teks tidak dapat diekstrak dari dokumen."
        }

    kategori_str = "\n".join(f"- {k}" for k in KATEGORI_INDUSTRI)
    prompt = f"""Kamu adalah sistem klasifikasi industri. Tugasmu: tentukan industri dari dokumen ini.

Nama file: {filename}

Kategori yang tersedia:
{kategori_str}

Cuplikan isi dokumen:
\"\"\"
{teks}
\"\"\"

Jawab HANYA dalam format JSON berikut, tanpa penjelasan tambahan:
{{
  "industri": "<pilih TEPAT satu dari kategori di atas>",
  "keyakinan": "<tinggi | sedang | rendah>",
  "alasan": "<penjelasan singkat 1 kalimat>"
}}"""

    try:
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=256,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = response.content[0].text.strip()

        # Parse JSON — bersihkan code block jika ada
        import json
        raw_clean = raw.replace("```json", "").replace("```", "").strip()
        return json.loads(raw_clean)

    except Exception as e:
        return {
            "industri": "ERROR",
            "keyakinan": "rendah",
            "alasan": f"API error: {str(e)[:80]}"
        }


def scan_dokumen(folder: Path) -> list[Path]:
    ekstensi = {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx"}
    if REKURSI_SUBFOLDER:
        files = [f for f in folder.rglob("*") if f.suffix.lower() in ekstensi]
    else:
        files = [f for f in folder.glob("*") if f.suffix.lower() in ekstensi]
    return sorted(files)


def main():
    # Validasi
    if not ANTHROPIC_API_KEY:
        print("[ERROR] ANTHROPIC_API_KEY belum diset.")
        print("Cara 1: set environment variable ANTHROPIC_API_KEY=sk-ant-...")
        print("Cara 2: isi langsung di baris ANTHROPIC_API_KEY di script ini.")
        return

    folder = Path(FOLDER_PATH)
    if not folder.exists():
        print(f"[ERROR] Folder tidak ditemukan: {FOLDER_PATH}")
        return

    # Scan file
    files = scan_dokumen(folder)
    print(f"✓ Ditemukan {len(files)} dokumen di: {FOLDER_PATH}\n")

    if not files:
        print("Tidak ada dokumen PDF/DOCX/XLSX/PPTX ditemukan.")
        return

    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

    # Buat/buka CSV output
    output_path = Path(OUTPUT_CSV)
    with open(output_path, "w", newline="", encoding="utf-8-sig") as csvfile:
        fieldnames = ["no", "nama_file", "ekstensi", "path_lengkap",
                      "industri", "keyakinan", "alasan", "status"]
        writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
        writer.writeheader()

        for i, filepath in enumerate(tqdm(files, desc="Memproses dokumen"), start=1):
            status = "OK"
            try:
                teks = ekstrak_teks(filepath)
                if not teks.strip():
                    status = "GAGAL_EKSTRAK"

                hasil = klasifikasi_industri(client, filepath.name, teks)
            except Exception:
                traceback.print_exc()
                hasil = {"industri": "ERROR", "keyakinan": "rendah", "alasan": "Exception saat proses"}
                status = "ERROR"

            writer.writerow({
                "no":            i,
                "nama_file":     filepath.name,
                "ekstensi":      filepath.suffix.lower(),
                "path_lengkap":  str(filepath),
                "industri":      hasil.get("industri", "?"),
                "keyakinan":     hasil.get("keyakinan", "?"),
                "alasan":        hasil.get("alasan", "?"),
                "status":        status,
            })
            csvfile.flush()  # Simpan tiap baris (aman jika script dihentikan)

            time.sleep(DELAY_BETWEEN_API)

    print(f"\n✅ Selesai! Hasil disimpan di: {output_path.resolve()}")
    print(f"   Total dokumen: {len(files)}")


if __name__ == "__main__":
    main()