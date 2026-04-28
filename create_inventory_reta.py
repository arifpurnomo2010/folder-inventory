import argparse
import asyncio
import json
import os
import re
import threading
import urllib.error
import urllib.request
from collections import Counter
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass
from html import unescape
from pathlib import Path
from queue import Queue, Empty
from typing import Callable

import pandas as pd


def load_env_file(env_path: Path = Path(".env")):
    if not env_path.exists():
        return

    for line in env_path.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        key = key.strip()
        value = value.strip().strip('"').strip("'")
        os.environ.setdefault(key, value)


load_env_file(Path(__file__).with_name(".env"))


# Configuration
DEFAULT_BASE_PATH = r'D:\Marketing Research\Reta'
DEFAULT_OUTPUT_FILE = r'D:\Marketing Research\Reta\Project Inventory.xlsx'

DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "").strip()
DEEPSEEK_API_URL = "https://api.deepseek.com/chat/completions"
DEEPSEEK_MODEL = "deepseek-chat"
MAX_CHARS_PER_DOC = 3000
MAX_CONCURRENT_REQUESTS = 20
API_TIMEOUT_SECONDS = 60
SUPPORTED_TEXT_EXTENSIONS = {
    '.pdf',
    '.docx',
    '.xlsx',
    '.xlsm',
    '.xls',
    '.xlsb',
    '.pptx',
    '.txt',
    '.csv',
    '.html',
    '.htm',
}

PRODUCT_CATEGORIES = [
    'Soft Drinks / Carbonated Drinks',
    'RTD Tea & Coffee',
    'Bottled Water / Mineral Water',
    'Juice & Nectar',
    'Energy & Isotonic Drinks',
    'Dairy Drinks',
    'Cooking Oil',
    'Condiments & Sauces',
    'Instant Noodles',
    'Biscuits & Cookies',
    'Confectionery / Candies',
    'Snacks',
    'Cereals & Breakfast',
    'Spices & Seasonings',
    'Liquid Milk',
    'Sweetened Condensed Milk',
    'Powdered Milk',
    'Yogurt & Cheese',
    'Baby Milk Formula',
    'Baby Food',
    'Diapers & Baby Care',
    'Skin Care',
    'Hair Care',
    'Oral Care',
    'Deodorant & Antiperspirant',
    'Feminine Hygiene',
    'Makeup',
    'Fragrance / Cologne',
    'Fabric Wash (Detergent)',
    'Dishwashing',
    'Floor & Surface Cleaner',
    'Air Freshener & Insecticide',
    'Health Supplements',
    'Analgesics & Cough',
    'Remedy Drinks',
    'Smartphones',
    'Feature Phones & Accessories',
    'Laptops / Notebooks',
    'Tablets',
    'Desktops & Monitors',
    'Printers & Peripherals',
    'Televisions',
    'Audio & Sound Systems',
    'Digital Camera & Imaging',
    'Refrigerators',
    'Washing Machines',
    'Air Conditioners',
    'Cooking Appliances',
    'Kitchen Appliances',
    'Home Care Appliances',
    'Personal Care Appliances',
    'Lighting',
    'Paint & Wall Covering',
    'Tools & Hardware',
    'Telecommunications Services',
    'Banking & Financial Services',
    'Insurance',
    'E-commerce / Online Retail',
    'Retail / Modern Trade',
    'Automotive (Vehicles)',
    'Travel / Aviation / Hospitality',
    'Real Estate / Property',
    'Energy / Oil & Gas / Utilities',
    'Logistics & Transportation',
    'Media & Entertainment',
    'Government / Public Sector',
    'Education',
    'Healthcare Services',
    'Pharmaceutical',
    'Manufacturing / Industrial',
    'Tobacco',
    'Research / Consulting',
    'Internal / Training',
    'Other',
]


def log_message(message: str, logger: Callable[[str], None] | None = None):
    if logger:
        logger(message)
    else:
        print(message)


@dataclass(frozen=True)
class TextExtractionResult:
    text: str
    ok: bool
    reason: str = ""


@dataclass(frozen=True)
class ProductCategoryResult:
    category: str
    extraction: TextExtractionResult


def extraction_success(text: str) -> TextExtractionResult:
    text = text[:MAX_CHARS_PER_DOC]
    if text.strip():
        return TextExtractionResult(text=text, ok=True, reason=f"{len(text)} chars extracted")
    return TextExtractionResult(text="", ok=False, reason="no extractable text")


def extraction_failure(reason: str) -> TextExtractionResult:
    return TextExtractionResult(text="", ok=False, reason=reason)


def check_required_dependencies(logger: Callable[[str], None] | None = None):
    missing = []
    checks = [
        ("pypdf", "pypdf"),
        ("docx", "python-docx"),
        ("openpyxl", "openpyxl"),
        ("pptx", "python-pptx"),
        ("xlrd", "xlrd"),
        ("pyxlsb", "pyxlsb"),
    ]

    for module_name, install_name in checks:
        try:
            __import__(module_name)
        except ImportError:
            missing.append(install_name)

    if missing:
        log_message("[ERROR] Missing libraries: " + ", ".join(missing), logger)
        log_message("Run: pip install pypdf python-docx openpyxl python-pptx pandas xlrd pyxlsb", logger)
        return False
    return True


def extract_text_pdf(path: Path) -> TextExtractionResult:
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
            if len(text) >= MAX_CHARS_PER_DOC:
                break
        result = extraction_success(text)
        if not result.ok:
            return extraction_failure("no extractable text (possibly scanned/image-only PDF)")
        return result
    except Exception as exc:
        return extraction_failure(f"PDF extraction failed: {str(exc)[:120]}")


def extract_text_docx(path: Path) -> TextExtractionResult:
    try:
        import docx

        doc = docx.Document(str(path))
        text = []
        for para in doc.paragraphs:
            text.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text.append(cell.text)
        return extraction_success("\n".join(text))
    except Exception as exc:
        return extraction_failure(f"DOCX extraction failed: {str(exc)[:120]}")


def extract_text_xlsx(path: Path) -> TextExtractionResult:
    try:
        import openpyxl

        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        text = ""
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            text += f"[Sheet: {sheet_name}]\n"
            for row in worksheet.iter_rows(values_only=True):
                line = " | ".join(str(cell) for cell in row if cell is not None)
                if line:
                    text += line + "\n"
                if len(text) >= MAX_CHARS_PER_DOC:
                    break
            if len(text) >= MAX_CHARS_PER_DOC:
                break
        workbook.close()
        return extraction_success(text)
    except Exception as exc:
        return extraction_failure(f"Excel extraction failed: {str(exc)[:120]}")


def extract_text_legacy_excel(path: Path) -> TextExtractionResult:
    try:
        sheets = pd.read_excel(path, sheet_name=None, header=None, nrows=200)
        text_parts = []
        for sheet_name, dataframe in sheets.items():
            text_parts.append(f"[Sheet: {sheet_name}]")
            for row in dataframe.fillna("").astype(str).itertuples(index=False, name=None):
                line = " | ".join(cell.strip() for cell in row if cell and cell.strip())
                if line:
                    text_parts.append(line)
                if len("\n".join(text_parts)) >= MAX_CHARS_PER_DOC:
                    break
            if len("\n".join(text_parts)) >= MAX_CHARS_PER_DOC:
                break
        return extraction_success("\n".join(text_parts))
    except Exception as exc:
        return extraction_failure(f"legacy Excel extraction failed: {str(exc)[:120]}")


def extract_text_plain(path: Path) -> TextExtractionResult:
    last_error = ""
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            text = path.read_text(encoding=encoding, errors="replace")
            return extraction_success(text)
        except Exception as exc:
            last_error = str(exc)[:120]
    return extraction_failure(f"text extraction failed: {last_error}")


def extract_text_html(path: Path) -> TextExtractionResult:
    result = extract_text_plain(path)
    if not result.ok:
        return result

    text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", result.text)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", unescape(text)).strip()
    return extraction_success(text)


def extract_text_pptx(path: Path) -> TextExtractionResult:
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        text = ""
        for slide_number, slide in enumerate(presentation.slides, start=1):
            text += f"[Slide {slide_number}]\n"
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        text += para.text + "\n"
                elif hasattr(shape, "text"):
                    text += shape.text + "\n"
            if len(text) >= MAX_CHARS_PER_DOC:
                break
        return extraction_success(text)
    except Exception as exc:
        return extraction_failure(f"PPTX extraction failed: {str(exc)[:120]}")


def extract_text(path: Path) -> TextExtractionResult:
    ext = path.suffix.lower()
    if ext == '.pdf':
        return extract_text_pdf(path)
    if ext == '.docx':
        return extract_text_docx(path)
    if ext in ('.xlsx', '.xlsm'):
        return extract_text_xlsx(path)
    if ext in ('.xls', '.xlsb'):
        return extract_text_legacy_excel(path)
    if ext == '.pptx':
        return extract_text_pptx(path)
    if ext in ('.txt', '.csv'):
        return extract_text_plain(path)
    if ext in ('.html', '.htm'):
        return extract_text_html(path)
    return extraction_failure(f"unsupported extension: {ext or '(none)'}")


def normalize_product_category(category: str) -> str:
    if not category:
        return 'Other'

    category_clean = category.strip()
    category_lookup = {item.lower(): item for item in PRODUCT_CATEGORIES}
    return category_lookup.get(category_clean.lower(), 'Other')


def parse_deepseek_json(raw_text: str) -> dict:
    raw_clean = raw_text.replace("```json", "").replace("```", "").strip()
    try:
        return json.loads(raw_clean)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", raw_clean, flags=re.DOTALL)
        if not match:
            raise
        return json.loads(match.group(0))


def classify_product_category_with_deepseek(
    filename: str,
    text: str,
    folder_context: str = "",
) -> str:
    has_text = bool(text.strip())
    has_context = bool(folder_context.strip())
    if not has_text and not has_context and not filename.strip():
        return 'Other'

    category_list = "\n".join(f"- {category}" for category in PRODUCT_CATEGORIES)
    text_block = text.strip() if has_text else "(tidak ada teks yang bisa diekstrak)"
    folder_block = folder_context.strip() if has_context else "(tidak diketahui)"
    prompt = f"""Kamu adalah sistem klasifikasi product/industry category untuk dokumen market research.

Tentukan kategori utama dokumen ini berdasarkan tiga sumber sinyal berikut, dengan urutan prioritas:
1. Cuplikan isi dokumen (jika tersedia)
2. Nama folder relatif (sering memuat nama klien/brand/topik)
3. Nama file

Path folder relatif: {folder_block}
Nama file: {filename}

Kategori yang tersedia:
{category_list}

Cuplikan isi dokumen:
\"\"\"
{text_block}
\"\"\"

Aturan penting:
- Pilih kategori paling relevan dari daftar jika ada sinyal industri/produk yang cukup dari sumber mana pun di atas.
- Jika konten kosong, klasifikasi tetap diusahakan dari path folder + nama file.
- Jangan memilih "Other" hanya karena confidence rendah; pilih kategori terdekat dan beri confidence rendah.
- Pilih "Other" hanya bila tidak ada sinyal yang cukup dari ketiga sumber tersebut.

Jawab HANYA dalam JSON valid, tanpa markdown dan tanpa penjelasan di luar JSON:
{{
  "product_category": "<pilih TEPAT satu dari kategori yang tersedia>",
  "confidence": "<tinggi | sedang | rendah>",
  "reason": "<alasan singkat 1 kalimat>"
}}"""

    payload = {
        "model": DEEPSEEK_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0,
        "max_tokens": 256,
        "response_format": {"type": "json_object"},
    }
    request = urllib.request.Request(
        DEEPSEEK_API_URL,
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
            "Content-Type": "application/json",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(request, timeout=API_TIMEOUT_SECONDS) as response:
            response_data = json.loads(response.read().decode("utf-8"))
        raw_content = response_data["choices"][0]["message"]["content"].strip()
        parsed = parse_deepseek_json(raw_content)
        return normalize_product_category(parsed.get("product_category", ""))
    except (urllib.error.URLError, urllib.error.HTTPError, KeyError, json.JSONDecodeError) as exc:
        print(f"[WARN] DeepSeek gagal untuk {filename}: {str(exc)[:120]}")
        return 'Other'

def get_file_format(filepath):
    ext = Path(filepath).suffix.lower()
    formats = {
        '.docx': 'File Word',
        '.doc': 'File Word',
        '.xlsx': 'File Excel',
        '.xls': 'File Excel',
        '.xlsm': 'File Excel (Macro)',
        '.xlsb': 'File Excel (Binary)',
        '.pptx': 'File PowerPoint',
        '.ppt': 'File PowerPoint',
        '.pdf': 'File PDF',
        '.txt': 'File Text',
        '.zip': 'File ZIP',
        '.rar': 'File RAR',
        '.jpeg': 'File Image',
        '.jpg': 'File Image',
        '.png': 'File Image',
        '.jfif': 'File Image',
        '.webp': 'File Image',
        '.svg': 'File Image',
        '.sav': 'File SPSS',
        '.csv': 'File CSV',
        '.html': 'File HTML',
        '.htm': 'File HTML',
        '.css': 'File CSS',
        '.download': 'File Download',
        '.tmp': 'File Temp'
    }
    return formats.get(ext, f'File {ext.upper()}' if ext else 'File')

def get_parent(path, base_path):
    rel_path = os.path.relpath(path, base_path)
    parts = Path(rel_path).parts
    if len(parts) <= 1:
        return 'ROOT'
    else:
        return parts[-2]


def build_folder_context(path: str, base_path: str) -> str:
    try:
        rel_path = os.path.relpath(path, base_path)
    except ValueError:
        rel_path = path
    parts = [part for part in Path(rel_path).parts if part not in ('', '.')]
    folder_parts = parts[:-1] if parts else []
    if not folder_parts:
        return 'ROOT'
    return ' / '.join(folder_parts)


def get_product_category_for_file(
    path: str,
    name: str,
    folder_context: str = "",
) -> ProductCategoryResult:
    ext = Path(path).suffix.lower()
    if ext not in SUPPORTED_TEXT_EXTENSIONS:
        return ProductCategoryResult(
            category='Other',
            extraction=extraction_failure(f"unsupported extension: {ext or '(none)'}"),
        )

    extraction = extract_text(Path(path))
    if not extraction.ok:
        return ProductCategoryResult(category='Other', extraction=extraction)

    category = classify_product_category_with_deepseek(name, extraction.text, folder_context)
    return ProductCategoryResult(category=category, extraction=extraction)


async def get_product_category_for_file_async(
    path: str,
    name: str,
    folder_context: str,
    semaphore: asyncio.Semaphore,
    executor: ThreadPoolExecutor,
    index: int,
    total: int,
    logger: Callable[[str], None] | None = None,
) -> str:
    ext = Path(path).suffix.lower()
    if ext not in SUPPORTED_TEXT_EXTENSIONS:
        return 'Other'

    async with semaphore:
        loop = asyncio.get_running_loop()
        result = await loop.run_in_executor(
            executor, get_product_category_for_file, path, name, folder_context
        )
        if result.extraction.ok:
            log_message(
                f"[{index}/{total}] {name[:60]} -> {result.category} ({result.extraction.reason})",
                logger,
            )
        else:
            log_message(
                f"[WARN] [{index}/{total}] {name[:60]} -> Other (read failed: {result.extraction.reason})",
                logger,
            )
        return result.category


async def classify_supported_files_async(
    all_items,
    base_path: str,
    max_concurrent: int = MAX_CONCURRENT_REQUESTS,
    logger: Callable[[str], None] | None = None,
):
    supported_files = [
        (path, os.path.basename(path), build_folder_context(path, base_path))
        for path, is_folder in all_items
        if not is_folder
        and not os.path.basename(path).startswith(('~$', '.'))
        and Path(path).suffix.lower() in SUPPORTED_TEXT_EXTENSIONS
    ]

    total = len(supported_files)
    if total == 0:
        log_message("No supported documents found for DeepSeek classification.", logger)
        return {}

    log_message(
        f"Classifying {total} supported documents with {max_concurrent} concurrent requests...",
        logger,
    )
    semaphore = asyncio.Semaphore(max_concurrent)
    categories = {}

    with ThreadPoolExecutor(max_workers=max_concurrent) as executor:
        tasks = [
            get_product_category_for_file_async(
                path, name, folder_context, semaphore, executor, index, total, logger
            )
            for index, (path, name, folder_context) in enumerate(supported_files, start=1)
        ]
        results = await asyncio.gather(*tasks)

    for (path, _name, _folder_context), category in zip(supported_files, results):
        categories[path] = category

    return categories


def build_inventory_dataframe(all_items, file_categories, base_path: str, logger: Callable[[str], None] | None = None):
    data = []
    id_counter = 1
    processed = 0

    log_message("Building inventory rows...", logger)
    for path, is_folder in all_items:
        name = os.path.basename(path)

        # Skip temporary and hidden files.
        if name.startswith('~$') or name.startswith('.'):
            continue

        if is_folder:
            entity_type = 'Folder'
            proposal = 'N'
            questionnaire = 'N'
            report = 'N'
            training = 'N'
            product_category = 'Other'
        else:
            entity_type = get_file_format(path)
            name_lower = name.lower()

            proposal = 'Y' if any(kw in name_lower for kw in ['proposal', 'prop']) else 'N'
            questionnaire = 'Y' if any(kw in name_lower for kw in ['questionnaire', 'qnr', "q're", 'guide']) else 'N'
            report = 'Y' if any(kw in name_lower for kw in ['report', 'topline', 'final report', 'presentation']) else 'N'
            training = 'Y' if any(kw in name_lower for kw in ['training', 'trainng', 'method', 'material']) else 'N'
            product_category = file_categories.get(path, 'Other')

        data.append({
            'ID': id_counter,
            'File/Folder_Name': name,
            'Entity_Name': name,
            'Entity_Type': entity_type,
            'Parent': get_parent(path, base_path),
            'Project_Name': name,
            'Proposal': proposal,
            'Questionnaire / Guide': questionnaire,
            'Report': report,
            'Method/Training material': training,
            'Location': path,
            'Product Category': product_category
        })

        id_counter += 1
        processed += 1
        if processed % 100 == 0:
            log_message(f"Processed {processed} items...", logger)

    return pd.DataFrame(data)


def inherit_from_siblings(
    df: pd.DataFrame,
    threshold: int = 2,
    logger: Callable[[str], None] | None = None,
) -> pd.DataFrame:
    log_message("Pass 2: inheriting categories from siblings...", logger)

    parent_to_categories: dict[str, list[str]] = {}
    for _, row in df.iterrows():
        if row['Entity_Type'] == 'Folder':
            continue
        category = row['Product Category']
        if not category or category == 'Other':
            continue
        parent_dir = os.path.dirname(row['Location'])
        parent_to_categories.setdefault(parent_dir, []).append(category)

    inherited = 0
    for idx, row in df.iterrows():
        if row['Entity_Type'] == 'Folder':
            continue
        if row['Product Category'] != 'Other':
            continue
        parent_dir = os.path.dirname(row['Location'])
        siblings = parent_to_categories.get(parent_dir, [])
        if not siblings:
            continue
        most_common, count = Counter(siblings).most_common(1)[0]
        if count >= threshold:
            df.at[idx, 'Product Category'] = most_common
            inherited += 1

    log_message(
        f"Pass 2: inherited categories for {inherited} files (threshold={threshold} siblings).",
        logger,
    )
    return df


async def classify_remaining_with_metadata_async(
    df: pd.DataFrame,
    base_path: str,
    max_concurrent: int = MAX_CONCURRENT_REQUESTS,
    logger: Callable[[str], None] | None = None,
) -> pd.DataFrame:
    targets = []
    for idx, row in df.iterrows():
        if row['Entity_Type'] == 'Folder':
            continue
        if row['Product Category'] != 'Other':
            continue
        name = row['File/Folder_Name']
        if not isinstance(name, str) or name.startswith('~$') or name.startswith('.'):
            continue
        targets.append((idx, row['Location'], name))

    total = len(targets)
    if total == 0:
        log_message("Pass 3: no remaining items needing metadata-only classification.", logger)
        return df

    log_message(
        f"Pass 3: classifying {total} remaining items via filename + folder path "
        f"with {max_concurrent} concurrent requests...",
        logger,
    )

    semaphore = asyncio.Semaphore(max_concurrent)

    async def classify_one(idx, path: str, name: str, index: int, executor: ThreadPoolExecutor):
        folder_context = build_folder_context(path, base_path)
        async with semaphore:
            loop = asyncio.get_running_loop()
            category = await loop.run_in_executor(
                executor,
                classify_product_category_with_deepseek,
                name,
                "",
                folder_context,
            )
            log_message(
                f"[meta {index}/{total}] {name[:60]} -> {category} "
                f"(folder={folder_context[:60]})",
                logger,
            )
            return idx, category

    with ThreadPoolExecutor(max_workers=max_concurrent) as executor:
        tasks = [
            classify_one(idx, path, name, i + 1, executor)
            for i, (idx, path, name) in enumerate(targets)
        ]
        results = await asyncio.gather(*tasks)

    rescued = 0
    for idx, category in results:
        df.at[idx, 'Product Category'] = category
        if category != 'Other':
            rescued += 1

    log_message(
        f"Pass 3: rescued {rescued} of {total} remaining items via metadata-only LLM.",
        logger,
    )
    return df


def update_folder_categories(
    df: pd.DataFrame,
    base_path: str,
    logger: Callable[[str], None] | None = None,
) -> pd.DataFrame:
    log_message("Updating folder categories based on contents...", logger)

    folder_paths = {}
    for idx, row in df.iterrows():
        if row['Entity_Type'] == 'Folder':
            folder_paths[row['Location']] = []

    for idx, row in df.iterrows():
        if row['Entity_Type'] != 'Folder':
            parent_path = os.path.dirname(row['Location'])
            while parent_path in folder_paths or parent_path.startswith(base_path):
                if parent_path in folder_paths:
                    folder_paths[parent_path].append(row['Product Category'])
                parent_path = os.path.dirname(parent_path)
                if parent_path == base_path:
                    break

    for idx, row in df.iterrows():
        if row['Entity_Type'] == 'Folder':
            loc = row['Location']
            if loc in folder_paths and folder_paths[loc]:
                known_categories = [
                    category
                    for category in folder_paths[loc]
                    if category and category != 'Other'
                ]
                category_pool = known_categories or folder_paths[loc]
                most_common = Counter(category_pool).most_common(1)[0][0]
                df.at[idx, 'Product Category'] = most_common

    return df


def scan_items(base_path: str, logger: Callable[[str], None] | None = None):
    log_message("Scanning all files and folders...", logger)
    all_items = []
    for root, dirs, files in os.walk(base_path):
        for folder_name in dirs:
            all_items.append((os.path.join(root, folder_name), True))
        for file_name in files:
            all_items.append((os.path.join(root, file_name), False))

    all_items.sort(key=lambda x: x[0])
    log_message(f"Found {len(all_items)} items", logger)
    return all_items


def run_inventory(
    base_path: str,
    output_file: str,
    max_concurrent: int = MAX_CONCURRENT_REQUESTS,
    logger: Callable[[str], None] | None = None,
):
    if not DEEPSEEK_API_KEY:
        log_message("[ERROR] DEEPSEEK_API_KEY is not set.", logger)
        log_message('PowerShell: $env:DEEPSEEK_API_KEY="<your_deepseek_api_key>"', logger)
        return False

    if not check_required_dependencies(logger):
        return False

    if not os.path.exists(base_path):
        log_message(f"[ERROR] Base folder not found: {base_path}", logger)
        return False

    output_parent = os.path.dirname(output_file)
    if output_parent and not os.path.exists(output_parent):
        log_message(f"[ERROR] Output folder not found: {output_parent}", logger)
        return False

    all_items = scan_items(base_path, logger)

    file_categories = asyncio.run(
        classify_supported_files_async(
            all_items,
            base_path=base_path,
            max_concurrent=max_concurrent,
            logger=logger,
        )
    )
    df = build_inventory_dataframe(all_items, file_categories, base_path, logger)

    df = inherit_from_siblings(df, logger=logger)

    df = asyncio.run(
        classify_remaining_with_metadata_async(
            df,
            base_path=base_path,
            max_concurrent=max_concurrent,
            logger=logger,
        )
    )

    df = update_folder_categories(df, base_path, logger)
    df.to_excel(output_file, index=False)
    log_message(f"Done! Created inventory with {len(df)} items", logger)
    log_message(f"Saved to: {output_file}", logger)
    log_message("Product Category distribution:", logger)
    log_message(str(df['Product Category'].value_counts()), logger)
    return True


def launch_gui():
    import tkinter as tk
    from tkinter import filedialog, messagebox

    root = tk.Tk()
    root.title("Folder Inventory Generator")
    root.geometry("760x520")

    log_queue = Queue()
    base_path_var = tk.StringVar(value=DEFAULT_BASE_PATH)
    output_file_var = tk.StringVar(value=DEFAULT_OUTPUT_FILE)
    concurrent_var = tk.StringVar(value=str(MAX_CONCURRENT_REQUESTS))
    running_var = tk.BooleanVar(value=False)

    def browse_base_path():
        selected = filedialog.askdirectory(title="Select Base Folder")
        if selected:
            base_path_var.set(selected)
            output_file_var.set(str(Path(selected) / "Project Inventory.xlsx"))

    def browse_output_file():
        selected = filedialog.asksaveasfilename(
            title="Select Output Excel File",
            defaultextension=".xlsx",
            filetypes=[("Excel Workbook", "*.xlsx")],
        )
        if selected:
            output_file_var.set(selected)

    def gui_logger(message: str):
        log_queue.put(message)

    def poll_log_queue():
        try:
            while True:
                message = log_queue.get_nowait()
                log_text.configure(state="normal")
                log_text.insert("end", message + "\n")
                log_text.see("end")
                log_text.configure(state="disabled")
        except Empty:
            pass
        root.after(100, poll_log_queue)

    def run_clicked():
        if running_var.get():
            return

        base_path = base_path_var.get().strip()
        output_file = output_file_var.get().strip()
        try:
            max_concurrent = int(concurrent_var.get().strip())
            if max_concurrent < 1:
                raise ValueError
        except ValueError:
            messagebox.showerror("Invalid Input", "Max concurrent requests must be a positive number.")
            return

        if not base_path or not output_file:
            messagebox.showerror("Invalid Input", "Base folder and output file are required.")
            return

        running_var.set(True)
        run_button.configure(state="disabled")
        log_text.configure(state="normal")
        log_text.delete("1.0", "end")
        log_text.configure(state="disabled")

        def worker():
            try:
                success = run_inventory(base_path, output_file, max_concurrent, gui_logger)
                gui_logger("Inventory generation completed." if success else "Inventory generation stopped.")
            except Exception as exc:
                gui_logger(f"[ERROR] Unexpected error: {exc}")
            finally:
                running_var.set(False)
                root.after(0, lambda: run_button.configure(state="normal"))

        threading.Thread(target=worker, daemon=True).start()

    form = tk.Frame(root, padx=16, pady=16)
    form.pack(fill="x")

    tk.Label(form, text="Base Folder").grid(row=0, column=0, sticky="w", pady=6)
    tk.Entry(form, textvariable=base_path_var, width=72).grid(row=0, column=1, sticky="ew", pady=6)
    tk.Button(form, text="Browse...", command=browse_base_path).grid(row=0, column=2, padx=(8, 0), pady=6)

    tk.Label(form, text="Output Excel File").grid(row=1, column=0, sticky="w", pady=6)
    tk.Entry(form, textvariable=output_file_var, width=72).grid(row=1, column=1, sticky="ew", pady=6)
    tk.Button(form, text="Save As...", command=browse_output_file).grid(row=1, column=2, padx=(8, 0), pady=6)

    tk.Label(form, text="Max Concurrent Requests").grid(row=2, column=0, sticky="w", pady=6)
    tk.Entry(form, textvariable=concurrent_var, width=10).grid(row=2, column=1, sticky="w", pady=6)

    run_button = tk.Button(form, text="Run Inventory", command=run_clicked, width=18)
    run_button.grid(row=3, column=1, sticky="w", pady=(12, 6))
    form.columnconfigure(1, weight=1)

    log_frame = tk.Frame(root, padx=16, pady=16)
    log_frame.pack(fill="both", expand=True)
    tk.Label(log_frame, text="Progress Log").pack(anchor="w")
    log_text = tk.Text(log_frame, height=18, state="disabled", wrap="word")
    log_text.pack(fill="both", expand=True)

    poll_log_queue()
    root.mainloop()


def parse_args():
    parser = argparse.ArgumentParser(description="Create a folder inventory with DeepSeek product categories.")
    parser.add_argument("--base-path", default=DEFAULT_BASE_PATH, help="Folder to scan.")
    parser.add_argument("--output-file", default=DEFAULT_OUTPUT_FILE, help="Excel output file path.")
    parser.add_argument(
        "--max-concurrent",
        type=int,
        default=MAX_CONCURRENT_REQUESTS,
        help="Maximum concurrent DeepSeek requests.",
    )
    parser.add_argument("--cli", action="store_true", help="Run without the graphical interface.")
    return parser.parse_args()


def main():
    args = parse_args()
    if args.cli:
        run_inventory(args.base_path, args.output_file, args.max_concurrent)
    else:
        launch_gui()


if __name__ == "__main__":
    main()