"""
classify_fast.py
────────────────
Klasifikasi industri dokumen lokal — VERSI CEPAT

MODE 1 — ASYNC (default): 20-50 request paralel → ~5-10 menit untuk 1000 dokumen
MODE 2 — BATCH API:        submit semua sekaligus → hasil dalam 1-24 jam, HARGA 50% lebih murah

SETUP:
    pip install anthropic pypdf python-docx openpyxl python-pptx tqdm

CARA PAKAI:
    Async (cepat):   python classify_fast.py --mode async
    Batch (murah):   python classify_fast.py --mode batch
    Cek hasil batch: python classify_fast.py --mode batch --check <BATCH_ID>
"""

import os
import csv
import json
import asyncio
import argparse
import traceback
from pathlib import Path
from datetime import datetime

# ── KONFIGURASI ───────────────────────────────────────────────────────────────
FOLDER_PATH       = r"C:\Users\NamaKamu\Documents\Dokumen"   # ← ganti ini
OUTPUT_CSV        = "hasil_klasifikasi.csv"
ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")

MAX_CHARS_PER_DOC = 2000   # Dikurangi dari 3000 — cukup untuk klasifikasi
MAX_CONCURRENT    = 30     # Jumlah request paralel (naikkan ke 50 jika tidak ada error)
REKURSI_SUBFOLDER = True

KATEGORI_INDUSTRI = [
    "Susu & Dairy", "Snack & Makanan Ringan", "Minuman", "Makanan & F&B",
    "Shampo & Perawatan Rambut", "Skincare & Kecantikan", "Kesehatan & Farmasi",
    "Otomotif (Motor)", "Otomotif (Mobil)", "Perbankan & Keuangan",
    "Asuransi", "Telekomunikasi", "E-commerce & Retail", "Properti",
    "Pendidikan", "Logistik & Transportasi", "Media & Hiburan",
    "Teknologi & Software", "Manufaktur", "Pemerintahan",
    "Lainnya"
]

# ─────────────────────────────────────────────────────────────────────────────

def ekstrak_teks(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            teks = ""
            for page in reader.pages:
                teks += page.extract_text() or ""
                if len(teks) >= MAX_CHARS_PER_DOC:
                    break
            return teks[:MAX_CHARS_PER_DOC]

        elif ext == ".docx":
            import docx
            doc = docx.Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)[:MAX_CHARS_PER_DOC]

        elif ext in (".xlsx", ".xlsm"):
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            teks = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    baris = " | ".join(str(c) for c in row if c is not None)
                    if baris:
                        teks += baris + "\n"
                    if len(teks) >= MAX_CHARS_PER_DOC:
                        break
                if len(teks) >= MAX_CHARS_PER_DOC:
                    break
            return teks[:MAX_CHARS_PER_DOC]

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            teks = ""
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            teks += para.text + "\n"
                if len(teks) >= MAX_CHARS_PER_DOC:
                    break
            return teks[:MAX_CHARS_PER_DOC]

    except Exception as e:
        return ""
    return ""


def buat_prompt(filename: str, teks: str) -> str:
    kategori_str = "\n".join(f"- {k}" for k in KATEGORI_INDUSTRI)
    isi = teks.strip() if teks.strip() else "(tidak dapat diekstrak)"
    return f"""Klasifikasikan industri dokumen ini. Jawab HANYA JSON, tanpa penjelasan.

Nama file: {filename}
Kategori tersedia:
{kategori_str}

Isi dokumen:
\"\"\"
{isi}
\"\"\"

Format respons (JSON saja):
{{"industri": "<pilih satu kategori>", "keyakinan": "<tinggi|sedang|rendah>", "alasan": "<1 kalimat>"}}"""


def parse_hasil(raw: str) -> dict:
    try:
        clean = raw.replace("```json", "").replace("```", "").strip()
        return json.loads(clean)
    except Exception:
        return {"industri": "ERROR", "keyakinan": "rendah", "alasan": f"Parse error: {raw[:60]}"}


def scan_dokumen(folder: Path) -> list:
    ekstensi = {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx"}
    if REKURSI_SUBFOLDER:
        return sorted(f for f in folder.rglob("*") if f.suffix.lower() in ekstensi)
    return sorted(f for f in folder.glob("*") if f.suffix.lower() in ekstensi)


# ── MODE 1: ASYNC ─────────────────────────────────────────────────────────────

async def proses_satu(client, semaphore, filepath: Path, index: int, total: int) -> dict:
    async with semaphore:
        teks = await asyncio.get_event_loop().run_in_executor(None, ekstrak_teks, filepath)
        prompt = buat_prompt(filepath.name, teks)
        status = "OK"

        try:
            response = await client.messages.create(
                model="claude-haiku-4-5-20251001",   # Haiku: 5x lebih cepat & murah vs Sonnet
                max_tokens=150,
                messages=[{"role": "user", "content": prompt}]
            )
            hasil = parse_hasil(response.content[0].text)
        except Exception as e:
            hasil = {"industri": "ERROR", "keyakinan": "rendah", "alasan": str(e)[:80]}
            status = "ERROR"

        print(f"  [{index}/{total}] {filepath.name[:50]} → {hasil.get('industri', '?')}")
        return {
            "no": index,
            "nama_file": filepath.name,
            "ekstensi": filepath.suffix.lower(),
            "path_lengkap": str(filepath),
            "industri": hasil.get("industri", "?"),
            "keyakinan": hasil.get("keyakinan", "?"),
            "alasan": hasil.get("alasan", "?"),
            "status": status if teks.strip() else "GAGAL_EKSTRAK",
        }


async def run_async(files: list):
    import anthropic

    client = anthropic.AsyncAnthropic(api_key=ANTHROPIC_API_KEY)
    semaphore = asyncio.Semaphore(MAX_CONCURRENT)
    total = len(files)

    print(f"🚀 Memproses {total} dokumen dengan {MAX_CONCURRENT} concurrent requests...\n")

    tasks = [proses_satu(client, semaphore, f, i+1, total) for i, f in enumerate(files)]
    hasil_list = await asyncio.gather(*tasks)

    # Urutkan berdasarkan nomor
    hasil_list = sorted(hasil_list, key=lambda x: x["no"])

    output_path = Path(OUTPUT_CSV)
    fieldnames = ["no", "nama_file", "ekstensi", "path_lengkap",
                  "industri", "keyakinan", "alasan", "status"]
    with open(output_path, "w", newline="", encoding="utf-8-sig") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(hasil_list)

    print(f"\n✅ Selesai! Hasil: {output_path.resolve()}")

    # Ringkasan per industri
    from collections import Counter
    counter = Counter(r["industri"] for r in hasil_list)
    print("\n📊 Ringkasan:")
    for industri, count in counter.most_common():
        print(f"   {count:>4}x  {industri}")


# ── MODE 2: BATCH API ─────────────────────────────────────────────────────────

def submit_batch(files: list):
    import anthropic

    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    print(f"📦 Menyiapkan {len(files)} request untuk Batch API...\n")

    requests = []
    for i, filepath in enumerate(files):
        teks = ekstrak_teks(filepath)
        prompt = buat_prompt(filepath.name, teks)
        requests.append({
            "custom_id": f"doc_{i}_{filepath.name[:40]}",
            "params": {
                "model": "claude-haiku-4-5-20251001",
                "max_tokens": 150,
                "messages": [{"role": "user", "content": prompt}]
            }
        })

    batch = client.messages.batches.create(requests=requests)
    batch_id = batch.id

    # Simpan metadata untuk keperluan polling
    meta = {
        "batch_id": batch_id,
        "submitted_at": datetime.now().isoformat(),
        "total_docs": len(files),
        "files": [{"index": i, "path": str(f)} for i, f in enumerate(files)]
    }
    meta_path = Path(f"batch_meta_{batch_id}.json")
    meta_path.write_text(json.dumps(meta, indent=2, ensure_ascii=False))

    print(f"✅ Batch submitted!")
    print(f"   Batch ID : {batch_id}")
    print(f"   Metadata : {meta_path}")
    print(f"\nCek hasil nanti dengan:")
    print(f"   python classify_fast.py --mode batch --check {batch_id}")


def check_batch(batch_id: str):
    import anthropic

    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    batch = client.messages.batches.retrieve(batch_id)
    status = batch.processing_status

    print(f"Batch {batch_id}: {status}")
    print(f"  Selesai   : {batch.request_counts.succeeded}")
    print(f"  Error     : {batch.request_counts.errored}")
    print(f"  Processing: {batch.request_counts.processing}")

    if status != "ended":
        print("\nBatch belum selesai. Coba lagi nanti.")
        return

    # Load metadata
    meta_path = Path(f"batch_meta_{batch_id}.json")
    if not meta_path.exists():
        print(f"[ERROR] File metadata tidak ditemukan: {meta_path}")
        return

    meta = json.loads(meta_path.read_text())
    file_map = {f["index"]: f["path"] for f in meta["files"]}

    results = []
    for result in client.messages.batches.results(batch_id):
        custom_id = result.custom_id
        index = int(custom_id.split("_")[1])
        filepath = Path(file_map.get(index, "unknown"))

        if result.result.type == "succeeded":
            raw = result.result.message.content[0].text
            hasil = parse_hasil(raw)
            status_row = "OK"
        else:
            hasil = {"industri": "ERROR", "keyakinan": "rendah", "alasan": "Batch error"}
            status_row = "ERROR"

        results.append({
            "no": index + 1,
            "nama_file": filepath.name,
            "ekstensi": filepath.suffix.lower(),
            "path_lengkap": str(filepath),
            "industri": hasil.get("industri", "?"),
            "keyakinan": hasil.get("keyakinan", "?"),
            "alasan": hasil.get("alasan", "?"),
            "status": status_row,
        })

    results.sort(key=lambda x: x["no"])

    output_path = Path(OUTPUT_CSV)
    fieldnames = ["no", "nama_file", "ekstensi", "path_lengkap",
                  "industri", "keyakinan", "alasan", "status"]
    with open(output_path, "w", newline="", encoding="utf-8-sig") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(results)

    print(f"\n✅ Hasil disimpan: {output_path.resolve()}")


# ── MAIN ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--mode", choices=["async", "batch"], default="async")
    parser.add_argument("--check", metavar="BATCH_ID", help="Cek & download hasil batch")
    args = parser.parse_args()

    if not ANTHROPIC_API_KEY:
        print("[ERROR] ANTHROPIC_API_KEY belum diset.")
        return

    if args.check:
        check_batch(args.check)
        return

    folder = Path(FOLDER_PATH)
    if not folder.exists():
        print(f"[ERROR] Folder tidak ditemukan: {FOLDER_PATH}")
        return

    files = scan_dokumen(folder)
    print(f"✓ Ditemukan {len(files)} dokumen\n")
    if not files:
        return

    if args.mode == "async":
        asyncio.run(run_async(files))
    else:
        submit_batch(files)


if __name__ == "__main__":
    main()